import streamlit as st
import os
import tempfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import google.generativeai as genai
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Image as RLImage, PageTemplate, Frame, BaseDocTemplate
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor
from reportlab.lib.utils import ImageReader
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from PIL import Image
from docx import Document
from docx.shared import Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
import io
import base64
from datetime import datetime
import re
import uuid
import logging
import traceback

# Configure Streamlit page
st.set_page_config(
    page_title="PPT to eBook Converter",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Configure logging (will be set dynamically based on debug mode)
logging.basicConfig(
    level=logging.INFO,  # Default to INFO level
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('ppt_to_ebook.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

def set_logging_level(debug_mode=False):
    """Dynamically set logging level based on debug mode"""
    level = logging.DEBUG if debug_mode else logging.INFO
    logging.getLogger().setLevel(level)
    for handler in logging.getLogger().handlers:
        handler.setLevel(level)
    logger.info(f"Logging level set to: {'DEBUG' if debug_mode else 'INFO'}")

# Custom CSS for modern UI
st.markdown("""
<style>
    .main-header {
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        padding: 2rem;
        border-radius: 10px;
        color: white;
        text-align: center;
        margin-bottom: 2rem;
    }
    
    .feature-card {
        background: #f8f9fa;
        padding: 1.5rem;
        border-radius: 10px;
        border-left: 4px solid #667eea;
        margin: 1rem 0;
    }
    
    .success-message {
        background: #d4edda;
        color: #155724;
        padding: 1rem;
        border-radius: 5px;
        border: 1px solid #c3e6cb;
    }
    
    .error-message {
        background: #f8d7da;
        color: #721c24;
        padding: 1rem;
        border-radius: 5px;
        border: 1px solid #f5c6cb;
    }
    
    .stButton > button {
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        color: white;
        border: none;
        padding: 0.5rem 2rem;
        border-radius: 25px;
        font-weight: bold;
    }
</style>
""", unsafe_allow_html=True)

class PPTToEBookConverter:
    def __init__(self, gemini_api_key):
        logger.info("Initializing PPTToEBookConverter")
        if gemini_api_key:
            try:
                logger.debug("Configuring Gemini API with provided key")
                genai.configure(api_key=gemini_api_key)
                self.model = genai.GenerativeModel('gemini-2.5-pro')
                logger.info("Gemini API configured successfully")
            except Exception as e:
                logger.error(f"Failed to configure Gemini API: {str(e)}")
                logger.debug(f"Gemini API error traceback: {traceback.format_exc()}")
                self.model = None
        else:
            logger.warning("No Gemini API key provided - AI enhancement will be disabled")
            self.model = None
    
    def extract_ppt_content(self, ppt_file):
        """Extract text content and images from PowerPoint slides"""
        logger.info("Starting PPT content extraction")
        try:
            logger.debug(f"Opening presentation file: {ppt_file.name if hasattr(ppt_file, 'name') else 'uploaded file'}")
            presentation = Presentation(ppt_file)
            slides_content = []
            total_slides = len(presentation.slides)
            logger.info(f"Found {total_slides} slides in presentation")
            
            for i, slide in enumerate(presentation.slides):
                logger.debug(f"Processing slide {i + 1}/{total_slides}")
                slide_content = {
                    'slide_number': i + 1,
                    'title': '',
                    'content': [],
                    'images': []
                }
                
                shape_count = len(slide.shapes)
                logger.debug(f"Slide {i + 1} has {shape_count} shapes")
                
                # Extract content from all shapes in the slide
                for j, shape in enumerate(slide.shapes):
                    logger.debug(f"Processing shape {j + 1}/{shape_count} on slide {i + 1}")
                    
                    # Handle text content
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        logger.debug(f"Found text content: {text[:50]}..." if len(text) > 50 else f"Found text content: {text}")
                        
                        if not slide_content['title'] and len(text) < 100:
                            slide_content['title'] = text
                            logger.debug(f"Set slide title: {text}")
                        else:
                            slide_content['content'].append(text)
                            logger.debug(f"Added content text (length: {len(text)})")
                    
                    # Handle images, diagrams, and flowcharts
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        logger.debug(f"Found image/diagram shape on slide {i + 1}")
                        try:
                            # Extract image data
                            image = shape.image
                            image_bytes = image.blob
                            image_size = len(image_bytes)
                            logger.debug(f"Extracted image data: {image_size} bytes")
                            
                            # Create a unique filename for the image
                            image_id = str(uuid.uuid4())
                            image_ext = image.ext if hasattr(image, 'ext') else 'png'
                            logger.debug(f"Image ID: {image_id}, Extension: {image_ext}")
                            
                            # Determine if it's likely a diagram/flowchart based on context
                            is_diagram = self._is_likely_diagram(shape, slide_content)
                            
                            # Store image data with slide title for proper captions
                            image_info = {
                                'id': image_id,
                                'data': image_bytes,
                                'extension': image_ext,
                                'filename': f"slide_{i+1}_{'diagram' if is_diagram else 'image'}_{image_id}.{image_ext}",
                                'type': 'diagram' if is_diagram else 'image',
                                'shape_name': shape.name if hasattr(shape, 'name') else None,
                                'slide_number': i + 1,  # Add slide number to image info
                                'slide_title': slide_content['title']  # Add slide title for captions
                            }
                            slide_content['images'].append(image_info)
                            logger.info(f"Successfully extracted {'diagram' if is_diagram else 'image'} from slide {i + 1}: {image_size} bytes, type: {image_bytes[:10] if image_bytes else 'None'}")
                            
                        except Exception as img_error:
                            logger.error(f"Failed to extract image/diagram from slide {i+1}: {str(img_error)}")
                            logger.debug(f"Image extraction error traceback: {traceback.format_exc()}")
                            st.warning(f"Could not extract image/diagram from slide {i+1}: {str(img_error)}")
                    
                    # Handle SmartArt diagrams and other graphic elements
                    elif hasattr(shape, 'shape_type') and shape.shape_type in [MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.AUTO_SHAPE]:
                        logger.debug(f"Found potential diagram/flowchart element on slide {i + 1}")
                        # Note: SmartArt and complex diagrams are harder to extract as images
                        # We'll log their presence for now
                        if hasattr(shape, 'name'):
                            logger.info(f"Found diagram element '{shape.name}' on slide {i + 1}")
                
                # If no title found, create one
                if not slide_content['title']:
                    slide_content['title'] = f"Slide {i + 1}"
                    logger.debug(f"Generated default title for slide {i + 1}")
                
                logger.info(f"Slide {i + 1} processed: {len(slide_content['content'])} text items, {len(slide_content['images'])} images")
                slides_content.append(slide_content)
            
            total_images = sum(len(slide.get('images', [])) for slide in slides_content)
            logger.info(f"PPT extraction completed: {len(slides_content)} slides, {total_images} total images")
            return slides_content
            
        except Exception as e:
            logger.error(f"Error extracting PPT content: {str(e)}")
            logger.debug(f"PPT extraction error traceback: {traceback.format_exc()}")
            st.error(f"Error extracting PPT content: {str(e)}")
            return None
    
    def _is_likely_diagram(self, shape, slide_content):
        """Determine if an image is likely a diagram or flowchart based on context"""
        try:
            # Check shape name for diagram-related keywords
            if hasattr(shape, 'name') and shape.name:
                diagram_keywords = ['diagram', 'flowchart', 'chart', 'flow', 'process', 'smartart', 'graphic']
                shape_name_lower = shape.name.lower()
                if any(keyword in shape_name_lower for keyword in diagram_keywords):
                    return True
            
            # Check slide content for diagram-related text
            all_text = slide_content.get('title', '') + ' ' + ' '.join(slide_content.get('content', []))
            diagram_text_keywords = ['diagram', 'flowchart', 'process', 'workflow', 'chart', 'flow', 'steps', 'procedure']
            if any(keyword in all_text.lower() for keyword in diagram_text_keywords):
                return True
            
            return False
        except Exception as e:
            logger.debug(f"Error determining diagram type: {str(e)}")
            return False
    
    def _clean_ai_response(self, response_text):
        """Clean up AI response by removing unwanted introductory phrases and fixing formatting"""
        if not response_text:
            return response_text
            
        # Common unwanted phrases to remove
        unwanted_phrases = [
            "Of course. Here is a comprehensive eBook chapter",
            "Here is a comprehensive eBook chapter",
            "Of course! Here is a comprehensive eBook chapter",
            "Here is a well-structured eBook chapter",
            "Of course. Here is a well-structured eBook chapter",
            "Here's a comprehensive eBook chapter",
            "Of course! Here's a comprehensive eBook chapter",
            "section based on the provided slide",
            "based on the provided slide title",
            "based on the slide content",
            "Here is the enhanced content",
            "Of course."
        ]
        
        cleaned_text = response_text.strip()
        
        # Remove unwanted phrases (case insensitive)
        for phrase in unwanted_phrases:
            # Remove at the beginning of text
            if cleaned_text.lower().startswith(phrase.lower()):
                cleaned_text = cleaned_text[len(phrase):].strip()
                # Remove any remaining colons or periods at the start
                cleaned_text = cleaned_text.lstrip(':. \n')
        
        # Remove any remaining introductory sentences that end with colon
        lines = cleaned_text.split('\n')
        if lines and ':' in lines[0] and len(lines[0]) < 100:
            # If first line is short and ends with colon, it's likely intro text
            if lines[0].strip().endswith(':'):
                cleaned_text = '\n'.join(lines[1:]).strip()
        
        # Fix common markdown formatting issues
        lines = cleaned_text.split('\n')
        fixed_lines = []
        
        for line in lines:
            line = line.strip()
            if line:
                # Fix headings that have inconsistent spacing
                if line.startswith('#'):
                    # Count the number of # symbols
                    hash_count = 0
                    for char in line:
                        if char == '#':
                            hash_count += 1
                        else:
                            break
                    
                    # Extract the heading text
                    heading_text = line[hash_count:].strip()
                    heading_text = heading_text.lstrip(': ').strip()
                    
                    # Reconstruct with proper spacing
                    if hash_count == 1:
                        fixed_lines.append(f"# {heading_text}")
                    elif hash_count == 2:
                        fixed_lines.append(f"## {heading_text}")
                    elif hash_count == 3:
                        fixed_lines.append(f"### {heading_text}")
                    elif hash_count >= 4:
                        fixed_lines.append(f"#### {heading_text}")
                else:
                    fixed_lines.append(line)
        
        cleaned_text = '\n'.join(fixed_lines)
        
        logger.debug(f"Cleaned AI response: removed intro text and fixed formatting, final length: {len(cleaned_text)}")
        return cleaned_text
    
    def _group_slides_into_chapters(self, slides_content, slides_per_chapter=5):
        """Group slides into chapters based on slides_per_chapter setting"""
        logger.info(f"Grouping {len(slides_content)} slides into chapters ({slides_per_chapter} slides per chapter)")
        
        chapters = []
        current_chapter = []
        chapter_number = 1
        
        for i, slide in enumerate(slides_content):
            current_chapter.append(slide)
            
            # Create a new chapter when we reach the limit or it's the last slide
            if len(current_chapter) >= slides_per_chapter or i == len(slides_content) - 1:
                # Combine content from all slides in this chapter
                combined_content = {
                    'slide_numbers': [s['slide_number'] for s in current_chapter],
                    'title': f"Chapter {chapter_number}",
                    'content': [],
                    'images': [],
                    'chapter_number': chapter_number
                }
                
                # Combine all slide content
                for slide in current_chapter:
                    # Add slide title as a section header
                    if slide['title'] and slide['title'] != f"Slide {slide['slide_number']}":
                        combined_content['content'].append(f"## {slide['title']}")
                    
                    # Add slide content
                    combined_content['content'].extend(slide['content'])
                    
                    # Add images with slide number preserved
                    for img in slide['images']:
                        img['slide_number'] = slide['slide_number']  # Ensure slide number is preserved
                        combined_content['images'].append(img)
                
                chapters.append(combined_content)
                logger.debug(f"Created chapter {chapter_number} with {len(current_chapter)} slides")
                
                # Reset for next chapter
                current_chapter = []
                chapter_number += 1
        
        logger.info(f"Created {len(chapters)} chapters from {len(slides_content)} slides")
        return chapters
    
    def _create_chapters_from_custom_ranges(self, slides_content, custom_ranges):
        """Create chapters based on custom slide ranges"""
        logger.info(f"Creating chapters from custom ranges: {custom_ranges}")
        
        chapters = []
        
        for chapter_num, (start_slide, end_slide) in enumerate(custom_ranges, 1):
            # Find slides in this range
            slides_in_range = []
            for slide in slides_content:
                slide_num = slide['slide_number']
                if start_slide <= slide_num <= end_slide:
                    slides_in_range.append(slide)
            
            if slides_in_range:
                # Create combined chapter content
                combined_content = {
                    'slide_numbers': [s['slide_number'] for s in slides_in_range],
                    'title': f"Chapter {chapter_num}",
                    'content': [],
                    'images': [],
                    'chapter_number': chapter_num
                }
                
                # Combine all slide content
                for slide in slides_in_range:
                    # Add slide title as a section header
                    if slide['title'] and slide['title'] != f"Slide {slide['slide_number']}":
                        combined_content['content'].append(f"## {slide['title']}")
                    
                    # Add slide content
                    combined_content['content'].extend(slide['content'])
                    
                    # Add images with slide number preserved
                    for img in slide['images']:
                        img['slide_number'] = slide['slide_number']  # Ensure slide number is preserved
                        combined_content['images'].append(img)
                
                chapters.append(combined_content)
                logger.debug(f"Created chapter {chapter_num} with slides {start_slide}-{end_slide} ({len(slides_in_range)} slides)")
            else:
                logger.warning(f"No slides found in range {start_slide}-{end_slide} for chapter {chapter_num}")
        
        logger.info(f"Created {len(chapters)} chapters from custom ranges")
        return chapters
    
    def enhance_content_with_gemini(self, slide_content, chapter_number=None):
        """Use Gemini API to enhance and structure slide content with image context"""
        # Handle both individual slides and grouped chapters
        if 'slide_numbers' in slide_content:
            # This is a grouped chapter
            slide_identifier = f"slides {min(slide_content['slide_numbers'])}-{max(slide_content['slide_numbers'])}"
            slide_num = slide_content['slide_numbers'][0]  # Use first slide number for logging
        else:
            # This is an individual slide
            slide_num = slide_content['slide_number']
            slide_identifier = f"slide {slide_num}"
            
        logger.info(f"Starting content enhancement for {slide_identifier}, chapter {chapter_number}")
        
        if not self.model:
            logger.warning(f"No Gemini model available for {slide_identifier} - using fallback content")
            # Return enhanced content structure even without AI
            if 'slide_numbers' in slide_content:
                # Grouped chapter fallback
                enhanced_content = {
                    'text': f"{slide_content.get('title', f'Chapter {chapter_number}')}\n\n" + '\n\n'.join(slide_content['content']),
                    'images': slide_content['images'],
                    'slide_numbers': slide_content['slide_numbers'],
                    'chapter_number': chapter_number
                }
            else:
                # Individual slide fallback
                enhanced_content = {
                    'text': f"{slide_content['title']}\n\n" + '\n\n'.join(slide_content['content']),
                    'images': slide_content['images'],
                    'slide_number': slide_content['slide_number'],
                    'chapter_number': chapter_number
                }
            logger.debug(f"Fallback content created for {slide_identifier}")
            return enhanced_content
        
        try:
            # Handle both individual slides and grouped chapters
            if 'slide_numbers' in slide_content:
                # This is a grouped chapter
                title = slide_content.get('title', f"Chapter {chapter_number}")
                content = '\n'.join(slide_content['content'])
                image_count = len(slide_content['images'])
            else:
                # This is an individual slide
                title = slide_content['title']
                content = '\n'.join(slide_content['content'])
                image_count = len(slide_content['images'])
            
            logger.debug(f"Slide {slide_num} - Title: {title[:50]}..., Content length: {len(content)}, Images: {image_count}")
            
            # Create context about images
            image_context = ""
            if image_count > 0:
                image_context = f"\n\nNote: This slide contains {image_count} image(s) that will be included in the eBook to support the content."
                logger.debug(f"Added image context for {image_count} images")
            
            # Create chapter-aware prompt
            chapter_info = f"This is Chapter {chapter_number}." if chapter_number else "This is a chapter section."
            
            prompt = f"""
            Transform this PowerPoint slide content into a comprehensive eBook chapter section.
            
            {chapter_info}
            
            Slide Title: {title}
            Content: {content}{image_context}
            
            Requirements:
            1. Create comprehensive, well-structured content suitable for an eBook
            2. If this is Chapter {chapter_number}, start with "Chapter {chapter_number}: [Title]" as the main heading (no # symbols)
            3. Use section headings like "[Title]" for major sections within the chapter (no ## symbols)
            4. Use subsection headings like "[Title]" for subsections (no ### symbols)
            5. Expand bullet points into full, coherent paragraphs
            6. Add context, explanations, and smooth transitions between ideas
            7. Maintain a professional, educational tone
            8. Structure the content logically for reading flow
            9. Do not repeat the same information multiple times
            10. Make the content engaging and informative
            11. Do not use markdown symbols (#, ##, ###) in headings - use clean text only
            
            {image_context}
            
            Generate clean, professional eBook content without any introductory phrases like "Here is" or "Of course".
            Format headings as clean text without markdown symbols.
            Start immediately with the chapter heading and content - no preamble.
            """
            
            logger.debug(f"Sending prompt to Gemini API for slide {slide_num} (prompt length: {len(prompt)})")
            response = self.model.generate_content(prompt)
            logger.info(f"Received Gemini API response for slide {slide_num} (response length: {len(response.text) if response.text else 0})")
            
            # Clean up the AI response
            cleaned_text = self._clean_ai_response(response.text)
            
            # Return enhanced content with images
            if 'slide_numbers' in slide_content:
                # Grouped chapter response
                enhanced_content = {
                    'text': cleaned_text,
                    'images': slide_content['images'],
                    'slide_numbers': slide_content['slide_numbers'],
                    'chapter_number': chapter_number
                }
                logger.debug(f"Enhanced grouped chapter content: {len(slide_content['images'])} images")
            else:
                # Individual slide response
                enhanced_content = {
                    'text': cleaned_text,
                    'images': slide_content['images'],
                    'slide_number': slide_content['slide_number'],
                    'chapter_number': chapter_number
                }
                logger.debug(f"Enhanced single slide content: {len(slide_content['images'])} images")
            
            logger.debug(f"Enhanced content structure created for {slide_identifier} with {len(enhanced_content['images'])} images")
            return enhanced_content
            
        except Exception as e:
            logger.error(f"Failed to enhance content with Gemini API for {slide_identifier}: {str(e)}")
            logger.debug(f"Gemini API error traceback for {slide_identifier}: {traceback.format_exc()}")
            st.warning(f"Could not enhance content with Gemini API: {str(e)}")
            
            # Fallback to original content with images
            if 'slide_numbers' in slide_content:
                # Grouped chapter error fallback
                enhanced_content = {
                    'text': f"{slide_content.get('title', f'Chapter {chapter_number}')}\n\n" + '\n\n'.join(slide_content['content']),
                    'images': slide_content['images'],
                    'slide_numbers': slide_content['slide_numbers'],
                    'chapter_number': chapter_number
                }
            else:
                # Individual slide error fallback
                enhanced_content = {
                    'text': f"{slide_content['title']}\n\n" + '\n\n'.join(slide_content['content']),
                    'images': slide_content['images'],
                    'slide_number': slide_content['slide_number'],
                    'chapter_number': chapter_number
                }
            logger.info(f"Using fallback content for {slide_identifier} due to API error")
            return enhanced_content
    
    def _create_custom_document(self, buffer, title, author, header_text, footer_text):
        """Create a custom PDF document with headers and footers"""
        logger.debug("Setting up custom document template with headers and footers")
        
        class CustomDocTemplate(BaseDocTemplate):
            def __init__(self, filename, **kwargs):
                BaseDocTemplate.__init__(self, filename, **kwargs)
                self.current_chapter = 1  # Track current chapter
                
                # Define frame for main content
                frame = Frame(
                    72, 72,  # x, y (left, bottom margins)
                    A4[0] - 144, A4[1] - 144,  # width, height (page size minus margins)
                    leftPadding=0, bottomPadding=0, rightPadding=0, topPadding=0
                )
                
                # Create page template with header and footer
                template = PageTemplate(
                    id='normal',
                    frames=[frame],
                    onPage=lambda canvas, doc: self._draw_header_footer(
                        canvas, doc, header_text, footer_text, title, author
                    )
                )
                
                self.addPageTemplates([template])
            
            def set_current_chapter(self, chapter_num):
                """Update the current chapter number"""
                self.current_chapter = chapter_num
            
            def _draw_header_footer(self, canvas, doc, header_text, footer_text, title, author):
                """Draw header and footer on each page"""
                canvas.saveState()
                
                # Header
                canvas.setFont('Helvetica-Bold', 10)
                canvas.setFillColor(HexColor('#2c3e50'))
                canvas.drawString(72, A4[1] - 50, header_text)
                
                # Header line
                canvas.setStrokeColor(HexColor('#bdc3c7'))
                canvas.setLineWidth(0.5)
                canvas.line(72, A4[1] - 55, A4[0] - 72, A4[1] - 55)
                
                # Footer
                canvas.setFont('Helvetica', 8)
                canvas.setFillColor(HexColor('#7f8c8d'))
                
                # Footer line
                canvas.line(72, 50, A4[0] - 72, 50)
                
                # Left footer (chapter number)
                chapter_text = f"Chapter {doc.current_chapter}"
                canvas.drawString(72, 35, chapter_text)
                
                # Right footer (page number)
                page_num = f"Page {doc.page}"
                canvas.drawRightString(A4[0] - 72, 35, page_num)
                
                canvas.restoreState()
        
        # Create the custom document
        doc = CustomDocTemplate(
            buffer,
            pagesize=A4,
            title=title,
            author=author
        )
        
        return doc
    
    class ChapterMarker:
        """Custom flowable to update chapter number in document"""
        def __init__(self, chapter_num):
            self.chapter_num = chapter_num
            
        def wrap(self, availWidth, availHeight):
            return (0, 0)  # Takes no space
            
        def draw(self):
            pass  # Nothing to draw
            
        def drawOn(self, canvas, x, y, _sW=0):
            # Update the document's current chapter
            if hasattr(canvas._doc, 'set_current_chapter'):
                canvas._doc.set_current_chapter(self.chapter_num)
        
        def getKeepWithNext(self):
            """Required method for ReportLab flowables"""
            return False
        
        def getSpaceAfter(self):
            """Required method for ReportLab flowables"""
            return 0
        
        def getSpaceBefore(self):
            """Required method for ReportLab flowables"""
            return 0
    
    def _parse_markdown_paragraph(self, paragraph, styles):
        """Parse a paragraph for Markdown formatting and return appropriate style and clean text"""
        paragraph = paragraph.strip()
        
        # Handle different heading levels - be more aggressive in cleaning
        if paragraph.startswith('#### ') or paragraph.startswith('####'):
            # H4 - Subsubsection
            clean_text = paragraph.lstrip('#').strip().replace('**', '').replace('*', '')
            # Remove any remaining colons or formatting
            clean_text = clean_text.lstrip(': ').strip()
            style = ParagraphStyle(
                'CustomH4',
                parent=styles['Heading4'],
                fontSize=14,
                spaceAfter=10,
                spaceBefore=15,
                textColor=HexColor('#2c3e50'),
                fontName='Helvetica-Bold'
            )
            return clean_text, style
            
        elif paragraph.startswith('### ') or paragraph.startswith('###'):
            # H3 - Subsection
            clean_text = paragraph.lstrip('#').strip().replace('**', '').replace('*', '')
            clean_text = clean_text.lstrip(': ').strip()
            style = ParagraphStyle(
                'CustomH3',
                parent=styles['Heading3'],
                fontSize=16,
                spaceAfter=12,
                spaceBefore=18,
                textColor=HexColor('#2c3e50'),
                fontName='Helvetica-Bold'
            )
            return clean_text, style
            
        elif paragraph.startswith('## ') or paragraph.startswith('##'):
            # H2 - Section
            clean_text = paragraph.lstrip('#').strip().replace('**', '').replace('*', '')
            clean_text = clean_text.lstrip(': ').strip()
            style = ParagraphStyle(
                'CustomH2',
                parent=styles['Heading2'],
                fontSize=18,
                spaceAfter=15,
                spaceBefore=20,
                textColor=HexColor('#2c3e50'),
                fontName='Helvetica-Bold'
            )
            return clean_text, style
            
        elif paragraph.startswith('# ') or paragraph.startswith('#'):
            # H1 - Chapter
            clean_text = paragraph.lstrip('#').strip().replace('**', '').replace('*', '')
            clean_text = clean_text.lstrip(': ').strip()
            style = ParagraphStyle(
                'CustomH1',
                parent=styles['Heading1'],
                fontSize=20,
                spaceAfter=15,
                spaceBefore=20,
                textColor=HexColor('#2c3e50'),
                fontName='Helvetica-Bold'
            )
            return clean_text, style
        
        else:
            # Regular paragraph - clean up any remaining markdown and stray # symbols
            clean_text = paragraph.replace('**', '').replace('*', '')
            # Remove any stray # symbols that might appear at the beginning
            clean_text = clean_text.lstrip('#').strip()
            clean_text = clean_text.lstrip(': ').strip()
            
            body_style = ParagraphStyle(
                'CustomBody',
                parent=styles['Normal'],
                fontSize=12,
                spaceAfter=12,
                leading=18,
                textColor=HexColor('#2c3e50'),
                alignment=0  # Left alignment
            )
            return clean_text, body_style
    
    def create_pdf_ebook(self, enhanced_content, title="Converted eBook", author="Unknown Author", header_text=None, footer_text=None):
        """Create a professional PDF eBook from enhanced content with images"""
        logger.info(f"Starting PDF eBook creation: '{title}'")
        logger.debug(f"Processing {len(enhanced_content)} content sections")
        
        buffer = io.BytesIO()
        
        # Set default header and footer if not provided
        if header_text is None:
            header_text = title
        if footer_text is None:
            footer_text = f"Generated by PPT to eBook Converter | Author: {author}"
        
        # Create custom PDF document with headers and footers
        logger.debug("Creating PDF document with custom template, headers, and footers")
        doc = self._create_custom_document(
            buffer, 
            title=title, 
            author=author, 
            header_text=header_text, 
            footer_text=footer_text
        )
        
        # Define styles
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            textColor=HexColor('#2c3e50'),
            alignment=1  # Center alignment
        )
        
        chapter_style = ParagraphStyle(
            'ChapterTitle',
            parent=styles['Heading2'],
            fontSize=18,
            spaceAfter=20,
            spaceBefore=20,
            textColor=HexColor('#34495e'),
            borderWidth=1,
            borderColor=HexColor('#bdc3c7'),
            borderPadding=10
        )
        
        body_style = ParagraphStyle(
            'CustomBody',
            parent=styles['Normal'],
            fontSize=12,
            spaceAfter=12,
            leading=18,
            textColor=HexColor('#2c3e50'),
            alignment=0  # Left alignment
        )
        
        image_caption_style = ParagraphStyle(
            'ImageCaption',
            parent=styles['Normal'],
            fontSize=10,
            spaceAfter=12,
            spaceBefore=6,
            textColor=HexColor('#7f8c8d'),
            alignment=1,  # Center alignment
            fontName='Helvetica-Oblique'
        )
        
        # Build PDF content
        story = []
        
        # Title page
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 0.3*inch))
        
        # Author information
        author_style = ParagraphStyle(
            'AuthorStyle',
            parent=styles['Normal'],
            fontSize=16,
            spaceAfter=20,
            textColor=HexColor('#34495e'),
            alignment=1  # Center alignment
        )
        story.append(Paragraph(f"by {author}", author_style))
        story.append(Spacer(1, 0.5*inch))
        
        # Generation date
        story.append(Paragraph(f"Generated on {datetime.now().strftime('%B %d, %Y')}", styles['Normal']))
        story.append(PageBreak())
        
        # Content sections
        for i, content in enumerate(enhanced_content):
            if i > 0:
                story.append(Spacer(1, 0.3*inch))
            
            # Handle both old string format and new dict format
            if isinstance(content, dict):
                content_text = content['text']
                content_images = content.get('images', [])
                slide_number = content.get('slide_number', i + 1)
                chapter_number = content.get('chapter_number', i + 1)
                logger.debug(f"PDF Section {i + 1}: Found {len(content_images)} images")
            else:
                # Fallback for old format
                content_text = content
                content_images = []
                slide_number = i + 1
                chapter_number = i + 1
                logger.debug(f"PDF Section {i + 1}: Fallback format, no images")
            
            # Add chapter marker to update footer
            story.append(self.ChapterMarker(chapter_number))
            
            # Extract title from content if it starts with a title-like pattern
            lines = content_text.split('\n')
            chapter_title = f"Chapter {i + 1}"
            content_body = content_text
            
            if lines and len(lines[0]) < 100 and not lines[0].startswith(' '):
                chapter_title = lines[0]
                content_body = '\n'.join(lines[1:])
            
            story.append(Paragraph(chapter_title, chapter_style))
            
            # Debug: Log content structure
            logger.debug(f"Content structure for section {i + 1}: {list(content.keys()) if isinstance(content, dict) else 'not dict'}")
            
            # For grouped chapters, process content with integrated images
            if 'slide_numbers' in content and isinstance(content.get('slide_numbers'), list):
                # This is a grouped chapter - process slide by slide with images
                logger.debug(f"Processing as grouped chapter with {len(content['slide_numbers'])} slides")
                self._process_grouped_chapter_content(story, content, styles, i + 1)
            else:
                # This is a single slide - process normally with images at the end
                logger.debug(f"Processing as single slide content")
                self._process_single_slide_content(story, content_text, content_images, styles, i + 1, slide_number)
            
            story.append(Spacer(1, 0.1*inch))
        
        # Build PDF
        logger.info(f"Building PDF document with {len(story)} elements")
        try:
            doc.build(story)
            buffer.seek(0)
            pdf_size = len(buffer.getvalue())
            logger.info(f"PDF eBook created successfully: {pdf_size} bytes")
            return buffer
        except Exception as e:
            logger.error(f"Failed to build PDF document: {str(e)}")
            logger.debug(f"PDF build error traceback: {traceback.format_exc()}")
            raise
    
    def create_docx_ebook(self, enhanced_content, title="Converted eBook", author="Unknown Author"):
        """Create a professional DOCX eBook from enhanced content with images"""
        logger.info(f"Starting DOCX eBook creation: '{title}'")
        logger.debug(f"Processing {len(enhanced_content)} content sections")
        
        # Create new document
        doc = Document()
        
        # Set document properties
        doc.core_properties.title = title
        doc.core_properties.author = author
        doc.core_properties.created = datetime.now()
        
        # Title page
        title_paragraph = doc.add_paragraph()
        title_run = title_paragraph.add_run(title)
        title_run.font.size = Inches(0.3)  # Large title
        title_run.bold = True
        title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Author
        author_paragraph = doc.add_paragraph()
        author_run = author_paragraph.add_run(f"by {author}")
        author_run.font.size = Inches(0.2)
        author_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Date
        date_paragraph = doc.add_paragraph()
        date_run = date_paragraph.add_run(f"Generated on {datetime.now().strftime('%B %d, %Y')}")
        date_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Page break
        doc.add_page_break()
        
        # Content sections
        for i, content in enumerate(enhanced_content):
            # Handle both old string format and new dict format
            if isinstance(content, dict):
                content_text = content['text']
                content_images = content.get('images', [])
                chapter_number = content.get('chapter_number', i + 1)
                logger.debug(f"DOCX Chapter {chapter_number}: Found {len(content_images)} images")
            else:
                # Fallback for old format
                content_text = content
                content_images = []
                chapter_number = i + 1
                logger.debug(f"DOCX Chapter {chapter_number}: Fallback format, no images")
            
            # Debug: Log content structure
            logger.debug(f"DOCX Content structure for chapter {chapter_number}: {list(content.keys()) if isinstance(content, dict) else 'not dict'}")
            
            # For grouped chapters, process content with integrated images
            if 'slide_numbers' in content and isinstance(content.get('slide_numbers'), list):
                # This is a grouped chapter - process slide by slide with images
                logger.debug(f"DOCX Processing as grouped chapter with {len(content['slide_numbers'])} slides")
                self._process_grouped_chapter_content_docx(doc, content, chapter_number)
            else:
                # This is a single slide - process normally with images at the end
                logger.debug(f"DOCX Processing as single slide content")
                self._process_single_slide_content_docx(doc, content_text, content_images, chapter_number)
        
        # Save to buffer
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        
        docx_size = len(buffer.getvalue())
        logger.info(f"DOCX eBook created successfully: {docx_size} bytes")
        return buffer
    
    def _process_grouped_chapter_content(self, story, content, styles, chapter_num):
        """Process grouped chapter content with images integrated per slide"""
        from reportlab.platypus import Spacer, Paragraph, Image as RLImage
        from reportlab.lib.units import inch
        from reportlab.lib.colors import HexColor
        from reportlab.lib.styles import ParagraphStyle
        
        logger.info(f"Processing grouped chapter {chapter_num} with {len(content.get('slide_numbers', []))} slides")
        
        # Get the original slide data to access slide titles and images
        slide_numbers = content.get('slide_numbers', [])
        all_images = content.get('images', [])
        enhanced_text = content.get('text', '')
        
        # Group images by their original slide number
        images_by_slide = {}
        slide_titles = {}  # Store slide titles for proper captions
        
        for img in all_images:
            slide_num = img.get('slide_number', 0)
            if slide_num not in images_by_slide:
                images_by_slide[slide_num] = []
            images_by_slide[slide_num].append(img)
            
            # Extract slide title from the original slide data if available
            if 'slide_title' in img:
                slide_titles[slide_num] = img['slide_title']
        
        # Split enhanced text into sections by slide headings
        text_sections = enhanced_text.split('## ')
        
        for i, section in enumerate(text_sections):
            if not section.strip():
                continue
                
            # Skip the first empty section if it exists
            if i == 0 and not section.strip():
                continue
                
            lines = section.strip().split('\n')
            if not lines:
                continue
                
            # First line is the slide heading
            slide_heading = lines[0].strip()
            
            # Try to extract slide number from heading or use sequential numbering
            current_slide_num = None
            if i <= len(slide_numbers):
                current_slide_num = slide_numbers[i-1] if i > 0 else slide_numbers[0]
            
            # Add the slide heading
            clean_text, paragraph_style = self._parse_markdown_paragraph(f"## {slide_heading}", styles)
            story.append(Paragraph(clean_text, paragraph_style))
            
            # Add images for this slide right after the heading
            if current_slide_num and current_slide_num in images_by_slide:
                logger.debug(f"Adding {len(images_by_slide[current_slide_num])} images for slide {current_slide_num}")
                # Pass slide title for proper captions
                slide_title = slide_titles.get(current_slide_num, slide_heading)
                self._add_images_to_story_with_title(story, images_by_slide[current_slide_num], chapter_num, current_slide_num, slide_title)
                # Remove processed images
                del images_by_slide[current_slide_num]
            
            # Add the rest of the content for this slide
            for line in lines[1:]:
                line = line.strip()
                if line and not (line.startswith('(Image:') or line.startswith('(Diagram:')):
                    clean_text, paragraph_style = self._parse_markdown_paragraph(line, styles)
                    story.append(Paragraph(clean_text, paragraph_style))
        
        # Add any remaining images that weren't processed
        for slide_num, images in images_by_slide.items():
            logger.debug(f"Adding {len(images)} remaining images from slide {slide_num}")
            slide_title = slide_titles.get(slide_num, f"Slide {slide_num}")
            self._add_images_to_story_with_title(story, images, chapter_num, slide_num, slide_title)
    
    def _process_single_slide_content(self, story, content_text, content_images, styles, chapter_num, slide_number):
        """Process single slide content with images at the end"""
        from reportlab.platypus import Spacer
        from reportlab.lib.units import inch
        
        logger.info(f"Processing single slide content for chapter {chapter_num}, slide {slide_number}")
        
        # Split content into paragraphs and process Markdown formatting
        paragraphs = [p.strip() for p in content_text.split('\n\n') if p.strip()]
        logger.debug(f"Processing {len(paragraphs)} paragraphs for chapter {chapter_num}")
        
        for paragraph in paragraphs:
            if paragraph:
                # Parse Markdown and get appropriate style
                clean_text, paragraph_style = self._parse_markdown_paragraph(paragraph, styles)
                logger.debug(f"Parsed paragraph: '{clean_text[:50]}...' with style: {paragraph_style.name}")
                story.append(Paragraph(clean_text, paragraph_style))
        
        # Add images at the end of the slide content
        if content_images:
            story.append(Spacer(1, 0.2*inch))
            self._add_images_to_story(story, content_images, chapter_num, slide_number)
    
    def _add_images_to_story(self, story, images, chapter_num, slide_number):
        """Add images to the PDF story with proper formatting"""
        from reportlab.platypus import Spacer, Paragraph, Image as RLImage
        from reportlab.lib.units import inch
        from reportlab.lib.colors import HexColor
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        
        styles = getSampleStyleSheet()
        
        logger.info(f"_add_images_to_story called: Adding {len(images)} images for chapter {chapter_num}, slide {slide_number}")
        
        if not images:
            logger.warning(f"No images to add for chapter {chapter_num}, slide {slide_number}")
            return
        
        for img_idx, image_data in enumerate(images):
            try:
                logger.debug(f"Processing image {img_idx + 1}/{len(images)} from slide {slide_number}")
                
                # Create PIL image from bytes
                img_bytes = io.BytesIO(image_data['data'])
                pil_image = Image.open(img_bytes)
                logger.debug(f"Opened PIL image: {pil_image.size} pixels, mode: {pil_image.mode}")
                
                # Calculate image dimensions for PDF
                img_width, img_height = pil_image.size
                aspect_ratio = img_width / img_height
                logger.debug(f"Original image size: {img_width}x{img_height}, aspect ratio: {aspect_ratio:.2f}")
                
                # Set maximum dimensions (fit within page margins)
                max_width = 400  # points (about 5.5 inches)
                max_height = 300  # points (about 4 inches)
                
                # Calculate new dimensions maintaining aspect ratio
                if img_width > max_width or img_height > max_height:
                    if aspect_ratio > 1:  # Wider than tall
                        new_width = max_width
                        new_height = max_width / aspect_ratio
                    else:  # Taller than wide
                        new_height = max_height
                        new_width = max_height * aspect_ratio
                else:
                    new_width = img_width
                    new_height = img_height
                
                logger.debug(f"Resized image to: {new_width:.1f}x{new_height:.1f} points")
                
                # Reset BytesIO for ReportLab
                img_bytes.seek(0)
                
                # Create ReportLab image
                rl_image = RLImage(img_bytes, width=new_width, height=new_height)
                story.append(rl_image)
                
                # Add image caption
                image_type = image_data.get('type', 'image')
                if image_type == 'diagram':
                    caption = f"Diagram {chapter_num}.{img_idx + 1}: Flowchart/Diagram from Slide {slide_number}"
                else:
                    caption = f"Figure {chapter_num}.{img_idx + 1}: Image from Slide {slide_number}"
                
                caption_style = ParagraphStyle(
                    'ImageCaption',
                    parent=styles['Normal'],
                    fontSize=9,
                    textColor=HexColor('#7f8c8d'),
                    alignment=1,  # Center alignment
                    spaceAfter=12
                )
                story.append(Paragraph(caption, caption_style))
                story.append(Spacer(1, 0.1*inch))
                logger.info(f"Successfully added image {img_idx + 1} from slide {slide_number} to PDF")
                
            except Exception as e:
                logger.error(f"Failed to process image {img_idx + 1} from slide {slide_number}: {str(e)}")
                logger.debug(f"Image processing error traceback: {traceback.format_exc()}")
                continue
    
    def _add_images_to_story_with_title(self, story, images, chapter_num, slide_number, slide_title):
        """Add images to the PDF story with proper slide title in captions"""
        from reportlab.platypus import Spacer, Paragraph, Image as RLImage
        from reportlab.lib.units import inch
        from reportlab.lib.colors import HexColor
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        
        styles = getSampleStyleSheet()
        
        logger.info(f"_add_images_to_story_with_title called: Adding {len(images)} images for chapter {chapter_num}, slide {slide_number}")
        
        if not images:
            logger.warning(f"No images to add for chapter {chapter_num}, slide {slide_number}")
            return
        
        for img_idx, image_data in enumerate(images):
            try:
                logger.debug(f"Processing image {img_idx + 1}/{len(images)} from slide {slide_number}")
                
                # Create PIL image from bytes
                img_bytes = io.BytesIO(image_data['data'])
                pil_image = Image.open(img_bytes)
                logger.debug(f"Opened PIL image: {pil_image.size} pixels, mode: {pil_image.mode}")
                
                # Calculate image dimensions for PDF
                img_width, img_height = pil_image.size
                aspect_ratio = img_width / img_height
                logger.debug(f"Original image size: {img_width}x{img_height}, aspect ratio: {aspect_ratio:.2f}")
                
                # Set maximum dimensions (fit within page margins)
                max_width = 400  # points (about 5.5 inches)
                max_height = 300  # points (about 4 inches)
                
                # Calculate new dimensions maintaining aspect ratio
                if img_width > max_width or img_height > max_height:
                    if aspect_ratio > 1:  # Wider than tall
                        new_width = max_width
                        new_height = max_width / aspect_ratio
                    else:  # Taller than wide
                        new_height = max_height
                        new_width = max_height * aspect_ratio
                else:
                    new_width = img_width
                    new_height = img_height
                
                logger.debug(f"Resized image to: {new_width:.1f}x{new_height:.1f} points")
                
                # Reset BytesIO for ReportLab
                img_bytes.seek(0)
                
                # Create ReportLab image
                rl_image = RLImage(img_bytes, width=new_width, height=new_height)
                story.append(rl_image)
                
                # Add image caption with proper slide title
                image_type = image_data.get('type', 'image')
                if image_type == 'diagram':
                    caption = f"Diagram {chapter_num}.{img_idx + 1}: {slide_title}"
                else:
                    caption = f"Figure {chapter_num}.{img_idx + 1}: {slide_title}"
                
                caption_style = ParagraphStyle(
                    'ImageCaption',
                    parent=styles['Normal'],
                    fontSize=9,
                    textColor=HexColor('#7f8c8d'),
                    alignment=1,  # Center alignment
                    spaceAfter=12
                )
                story.append(Paragraph(caption, caption_style))
                story.append(Spacer(1, 0.1*inch))
                logger.info(f"Successfully added image {img_idx + 1} from slide {slide_number} to PDF with title: {slide_title}")
                
            except Exception as e:
                logger.error(f"Failed to process image {img_idx + 1} from slide {slide_number}: {str(e)}")
                logger.debug(f"Image processing error traceback: {traceback.format_exc()}")
                continue
    
    def _process_grouped_chapter_content_docx(self, doc, content, chapter_num):
        """Process grouped chapter content for DOCX with images integrated per slide"""
        logger.info(f"Processing grouped chapter {chapter_num} for DOCX with {len(content.get('slide_numbers', []))} slides")
        
        # Get the original slide data to access slide titles and images
        slide_numbers = content.get('slide_numbers', [])
        all_images = content.get('images', [])
        enhanced_text = content.get('text', '')
        
        # Group images by their original slide number
        images_by_slide = {}
        slide_titles = {}  # Store slide titles for proper captions
        
        for img in all_images:
            slide_num = img.get('slide_number', 0)
            if slide_num not in images_by_slide:
                images_by_slide[slide_num] = []
            images_by_slide[slide_num].append(img)
            
            # Extract slide title from the original slide data if available
            if 'slide_title' in img:
                slide_titles[slide_num] = img['slide_title']
        
        # Split enhanced text into sections by slide headings
        text_sections = enhanced_text.split('## ')
        
        for i, section in enumerate(text_sections):
            if not section.strip():
                continue
                
            # Skip the first empty section if it exists
            if i == 0 and not section.strip():
                continue
                
            lines = section.strip().split('\n')
            if not lines:
                continue
                
            # First line is the slide heading
            slide_heading = lines[0].strip()
            
            # Try to extract slide number from heading or use sequential numbering
            current_slide_num = None
            if i <= len(slide_numbers):
                current_slide_num = slide_numbers[i-1] if i > 0 else slide_numbers[0]
            
            # Add the slide heading
            self._add_docx_paragraph(doc, f"## {slide_heading}")
            
            # Add images for this slide right after the heading
            if current_slide_num and current_slide_num in images_by_slide:
                logger.debug(f"Adding {len(images_by_slide[current_slide_num])} images for slide {current_slide_num}")
                # Pass slide title for proper captions
                slide_title = slide_titles.get(current_slide_num, slide_heading)
                self._add_images_to_docx_with_title(doc, images_by_slide[current_slide_num], chapter_num, current_slide_num, slide_title)
                # Remove processed images
                del images_by_slide[current_slide_num]
            
            # Add the rest of the content for this slide
            for line in lines[1:]:
                line = line.strip()
                if line and not (line.startswith('(Image:') or line.startswith('(Diagram:')):
                    self._add_docx_paragraph(doc, line)
        
        # Add any remaining images that weren't processed
        for slide_num, images in images_by_slide.items():
            logger.debug(f"Adding {len(images)} remaining images from slide {slide_num}")
            slide_title = slide_titles.get(slide_num, f"Slide {slide_num}")
            self._add_images_to_docx_with_title(doc, images, chapter_num, slide_num, slide_title)
    
    def _process_single_slide_content_docx(self, doc, content_text, content_images, chapter_num):
        """Process single slide content for DOCX with images at the end"""
        logger.info(f"Processing single slide content for DOCX chapter {chapter_num}")
        
        # Split content into paragraphs and process Markdown formatting
        paragraphs = [p.strip() for p in content_text.split('\n\n') if p.strip()]
        logger.debug(f"Processing {len(paragraphs)} paragraphs for DOCX chapter {chapter_num}")
        
        for paragraph_text in paragraphs:
            if paragraph_text:
                self._add_docx_paragraph(doc, paragraph_text)
        
        # Add images at the end of the slide content
        if content_images:
            # Get slide number from first image or use chapter number
            slide_number = content_images[0].get('slide_number', chapter_num) if content_images else chapter_num
            self._add_images_to_docx(doc, content_images, chapter_num, slide_number)
    
    def _add_docx_paragraph(self, doc, paragraph_text):
        """Add a paragraph to DOCX with proper Markdown formatting"""
        clean_text = paragraph_text.strip()
        
        # Remove Markdown symbols and apply formatting - be more aggressive in cleaning
        if clean_text.startswith('#### ') or clean_text.startswith('####'):
            clean_text = clean_text.lstrip('#').strip().replace('**', '').replace('*', '')
            clean_text = clean_text.lstrip(': ').strip()
            heading_paragraph = doc.add_paragraph()
            heading_run = heading_paragraph.add_run(clean_text)
            heading_run.bold = True
            heading_run.font.size = Inches(0.15)
        elif clean_text.startswith('### ') or clean_text.startswith('###'):
            clean_text = clean_text.lstrip('#').strip().replace('**', '').replace('*', '')
            clean_text = clean_text.lstrip(': ').strip()
            heading_paragraph = doc.add_paragraph()
            heading_run = heading_paragraph.add_run(clean_text)
            heading_run.bold = True
            heading_run.font.size = Inches(0.17)
        elif clean_text.startswith('## ') or clean_text.startswith('##'):
            clean_text = clean_text.lstrip('#').strip().replace('**', '').replace('*', '')
            clean_text = clean_text.lstrip(': ').strip()
            heading_paragraph = doc.add_paragraph()
            heading_run = heading_paragraph.add_run(clean_text)
            heading_run.bold = True
            heading_run.font.size = Inches(0.2)
        elif clean_text.startswith('# ') or clean_text.startswith('#'):
            clean_text = clean_text.lstrip('#').strip().replace('**', '').replace('*', '')
            clean_text = clean_text.lstrip(': ').strip()
            heading_paragraph = doc.add_paragraph()
            heading_run = heading_paragraph.add_run(clean_text)
            heading_run.bold = True
            heading_run.font.size = Inches(0.25)
        else:
            # Regular paragraph - clean up any remaining markdown and stray # symbols
            clean_text = clean_text.replace('**', '').replace('*', '')
            # Remove any stray # symbols that might appear at the beginning
            clean_text = clean_text.lstrip('#').strip()
            clean_text = clean_text.lstrip(': ').strip()
            doc.add_paragraph(clean_text)
    
    def _add_images_to_docx(self, doc, images, chapter_num, slide_number):
        """Add images to DOCX document with proper formatting"""
        logger.info(f"_add_images_to_docx called: Adding {len(images)} images to DOCX for chapter {chapter_num}, slide {slide_number}")
        
        if not images:
            logger.warning(f"No images to add to DOCX for chapter {chapter_num}, slide {slide_number}")
            return
        
        for img_idx, image_data in enumerate(images):
            try:
                img_bytes = io.BytesIO(image_data['data'])
                
                # Add image to document
                paragraph = doc.add_paragraph()
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                
                # Add image with reasonable size
                run.add_picture(img_bytes, width=Inches(5))
                
                # Add caption
                image_type = image_data.get('type', 'image')
                if image_type == 'diagram':
                    caption = f"Diagram {chapter_num}.{img_idx + 1}: Flowchart/Diagram from Slide {slide_number}"
                else:
                    caption = f"Figure {chapter_num}.{img_idx + 1}: Image from Slide {slide_number}"
                
                caption_paragraph = doc.add_paragraph(caption)
                caption_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                caption_run = caption_paragraph.runs[0]
                caption_run.italic = True
                
                logger.info(f"Successfully added image {img_idx + 1} from slide {slide_number} to DOCX")
                
            except Exception as e:
                logger.error(f"Failed to add image {img_idx + 1} from slide {slide_number} to DOCX: {str(e)}")
                continue
    
    def _add_images_to_docx_with_title(self, doc, images, chapter_num, slide_number, slide_title):
        """Add images to DOCX document with proper slide title in captions"""
        logger.info(f"_add_images_to_docx_with_title called: Adding {len(images)} images to DOCX for chapter {chapter_num}, slide {slide_number}")
        
        if not images:
            logger.warning(f"No images to add to DOCX for chapter {chapter_num}, slide {slide_number}")
            return
        
        for img_idx, image_data in enumerate(images):
            try:
                img_bytes = io.BytesIO(image_data['data'])
                
                # Add image to document
                paragraph = doc.add_paragraph()
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                
                # Add image with reasonable size
                run.add_picture(img_bytes, width=Inches(5))
                
                # Add caption with proper slide title
                image_type = image_data.get('type', 'image')
                if image_type == 'diagram':
                    caption = f"Diagram {chapter_num}.{img_idx + 1}: {slide_title}"
                else:
                    caption = f"Figure {chapter_num}.{img_idx + 1}: {slide_title}"
                
                caption_paragraph = doc.add_paragraph(caption)
                caption_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                caption_run = caption_paragraph.runs[0]
                caption_run.italic = True
                
                logger.info(f"Successfully added image {img_idx + 1} from slide {slide_number} to DOCX with title: {slide_title}")
                
            except Exception as e:
                logger.error(f"Failed to add image {img_idx + 1} from slide {slide_number} to DOCX: {str(e)}")
                continue

def main():
    # Header
    st.markdown("""
    <div class="main-header">
        <h1>📚 PPT to eBook Converter</h1>
        <p>Transform your PowerPoint presentations into professional PDF eBooks using AI</p>
    </div>
    """, unsafe_allow_html=True)
    
    # Sidebar for configuration
    with st.sidebar:
        st.header("⚙️ Configuration")
        
        # Gemini API Key input
        gemini_api_key = st.text_input(
            "Gemini API Key",
            type="password",
            help="Enter your Google Gemini API key to enhance content with AI"
        )
        
        if not gemini_api_key:
            st.warning("⚠️ Add Gemini API key to enable AI content enhancement")
        
        st.markdown("---")
        
        # eBook settings
        st.subheader("📚 eBook Settings")
        ebook_title = st.text_input("eBook Title", value="My eBook")
        ebook_author = st.text_input("Author Name", value="Anonymous")
        
        # Output format selection
        output_format = st.radio(
            "Output Format",
            ["PDF", "DOCX", "Both"],
            help="Choose the format for your eBook output"
        )
        
        # Advanced settings
        with st.expander("⚙️ Advanced Settings"):
            custom_header = st.text_input(
                "Custom Header Text", 
                placeholder="Leave empty to use eBook title",
                help="Text to appear at the top of each page"
            )
            custom_footer = st.text_input(
                "Custom Footer Text", 
                placeholder="Leave empty for default footer",
                help="Custom text for footer (page numbers will still appear)"
            )
            
            st.markdown("**Chapter Organization**")
            
            # Chapter organization method
            chapter_method = st.radio(
                "Chapter Organization Method",
                ["Automatic (Equal Groups)", "Custom Ranges", "One Slide Per Chapter"],
                help="Choose how to organize slides into chapters"
            )
            
            if chapter_method == "Automatic (Equal Groups)":
                slides_per_chapter = st.slider(
                    "Slides per Chapter",
                    min_value=2,
                    max_value=20,
                    value=5,
                    help="Group multiple slides into single chapters. Useful for large presentations."
                )
                st.info(f"📚 With {slides_per_chapter} slides per chapter, a 150-slide presentation will create ~{150//slides_per_chapter + (1 if 150%slides_per_chapter else 0)} chapters.")
                custom_ranges = None
                
            elif chapter_method == "Custom Ranges":
                st.markdown("**Define Custom Chapter Ranges:**")
                st.caption("Format: 'start-end' for each chapter, separated by commas")
                st.caption("Example: '1-5, 6-12, 13-18' creates 3 chapters")
                
                custom_ranges_input = st.text_area(
                    "Chapter Ranges",
                    placeholder="1-5, 6-12, 13-18, 19-25",
                    help="Define slide ranges for each chapter. Example: '1-5, 6-12' creates Chapter 1 (slides 1-5) and Chapter 2 (slides 6-12)"
                )
                
                if custom_ranges_input.strip():
                    try:
                        # Parse custom ranges
                        ranges = [r.strip() for r in custom_ranges_input.split(',')]
                        custom_ranges = []
                        for r in ranges:
                            if '-' in r:
                                start, end = map(int, r.split('-'))
                                custom_ranges.append((start, end))
                        
                        st.success(f"✅ Will create {len(custom_ranges)} chapters with custom ranges")
                        for i, (start, end) in enumerate(custom_ranges, 1):
                            st.caption(f"Chapter {i}: Slides {start}-{end} ({end-start+1} slides)")
                    except Exception as e:
                        st.error(f"❌ Invalid format: {str(e)}. Use format like '1-5, 6-12'")
                        custom_ranges = None
                else:
                    custom_ranges = None
                    st.info("📝 Enter chapter ranges above to see preview")
                
                slides_per_chapter = None  # Not used for custom ranges
                
            else:  # One Slide Per Chapter
                slides_per_chapter = 1
                custom_ranges = None
                st.info("📚 Each slide will become a separate chapter")
        
        st.markdown("---")
        
        # Features
        st.markdown("""
        ### ✨ Features
        - 📤 Upload PPTX files (up to 100MB)
        - 🤖 AI-powered content enhancement
        - 📸 Images & diagrams extraction
        - 📈 Flowchart & diagram recognition
        - 📄 Multiple output formats (PDF & DOCX)
        - 📚 Flexible chapter organization
        - 👤 Author information & custom branding
        - 📝 Chapter tracking in footers
        - 🎨 Modern eBook formatting
        - 📱 Responsive design
        """)
        
        st.markdown("---")
        
        # Debug section
        debug_mode = st.checkbox("🔍 Debug Mode", help="Enable detailed logging for troubleshooting")
        
        # Set logging level based on debug mode
        set_logging_level(debug_mode)
        
        if debug_mode:
            st.subheader("📝 Recent Debug Logs")
            st.info("ℹ️ Debug mode enabled - detailed logs will be captured during processing")
            
            try:
                if os.path.exists('ppt_to_ebook.log'):
                    with open('ppt_to_ebook.log', 'r') as f:
                        logs = f.readlines()
                        # Show last 30 lines in debug mode
                        recent_logs = logs[-30:] if len(logs) > 30 else logs
                        log_text = ''.join(recent_logs)
                        st.text_area("Debug Log Output", log_text, height=300, disabled=True)
                        
                        # Show log statistics
                        total_lines = len(logs)
                        debug_lines = len([line for line in logs if 'DEBUG' in line])
                        info_lines = len([line for line in logs if 'INFO' in line])
                        error_lines = len([line for line in logs if 'ERROR' in line])
                        
                        col1, col2, col3, col4 = st.columns(4)
                        with col1:
                            st.metric("Total Logs", total_lines)
                        with col2:
                            st.metric("Debug", debug_lines)
                        with col3:
                            st.metric("Info", info_lines)
                        with col4:
                            st.metric("Errors", error_lines)
                else:
                    st.info("No log file found yet. Start processing to see detailed logs.")
            except Exception as e:
                st.error(f"Could not read log file: {str(e)}")
        else:
            st.caption("📊 Normal mode - only important messages are logged")
    
    # Main content area
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.markdown("""
        <div class="feature-card">
            <h3>🚀 How it works</h3>
            <ol>
                <li>Upload your PowerPoint presentation</li>
                <li>Extract text content and images from slides</li>
                <li>AI enhances content for eBook format</li>
                <li>Generate professional PDF with images</li>
                <li>Download your complete eBook!</li>
            </ol>
        </div>
        """, unsafe_allow_html=True)
        
        # File upload
        uploaded_file = st.file_uploader(
            "Choose a PowerPoint file",
            type=['pptx'],
            help="Upload a PPTX file (maximum 100MB)"
        )
        
        if uploaded_file is not None:
            # Check file size
            file_size = len(uploaded_file.getvalue())
            if file_size > 100 * 1024 * 1024:  # 100MB limit
                st.error("❌ File size exceeds 100MB limit. Please upload a smaller file.")
                return
            
            st.success(f"✅ File uploaded successfully! Size: {file_size / (1024*1024):.2f} MB")
            
            # Process button
            if st.button("🔄 Convert to eBook", key="convert_btn"):
                logger.info(f"Starting eBook conversion process for file: {uploaded_file.name}")
                with st.spinner("Processing your presentation..."):
                    # Initialize converter
                    logger.debug("Initializing PPTToEBookConverter")
                    converter = PPTToEBookConverter(gemini_api_key)
                    
                    # Extract content from PPT
                    st.info("📚 Extracting content from slides...")
                    logger.info("Starting slide content extraction")
                    slides_content = converter.extract_ppt_content(uploaded_file)
                    
                    if slides_content:
                        # Count total images extracted
                        total_images = sum(len(slide.get('images', [])) for slide in slides_content)
                        logger.info(f"Extraction completed: {len(slides_content)} slides, {total_images} images")
                        
                        # Debug: Log images per slide
                        for idx, slide in enumerate(slides_content):
                            slide_images = len(slide.get('images', []))
                            if slide_images > 0:
                                logger.debug(f"Slide {idx + 1}: {slide_images} images extracted")
                        
                        st.success(f"✅ Extracted content from {len(slides_content)} slides with {total_images} images")
                        
                        # Show extraction summary
                        if total_images > 0:
                            st.info(f"📸 Found {total_images} images that will be included in your eBook")
                        
                        # Organize slides into chapters based on user preference
                        if chapter_method == "Custom Ranges" and custom_ranges:
                            logger.info(f"Creating chapters from custom ranges: {custom_ranges}")
                            chapters = converter._create_chapters_from_custom_ranges(slides_content, custom_ranges)
                            st.info(f"📚 Organized into {len(chapters)} chapters using custom ranges")
                        elif chapter_method == "One Slide Per Chapter":
                            logger.info("Creating individual chapters for each slide")
                            # Convert each slide to individual chapter format
                            chapters = []
                            for i, slide in enumerate(slides_content, 1):
                                chapter = {
                                    'slide_numbers': [slide['slide_number']],
                                    'title': slide['title'],
                                    'content': slide['content'],
                                    'images': slide['images'],
                                    'chapter_number': i
                                }
                                chapters.append(chapter)
                            st.info(f"📚 Created {len(chapters)} individual chapters (one per slide)")
                        else:
                            # Automatic equal groups
                            logger.info(f"Grouping slides into chapters: {slides_per_chapter} slides per chapter")
                            chapters = converter._group_slides_into_chapters(slides_content, slides_per_chapter)
                            st.info(f"📚 Organized into {len(chapters)} chapters ({slides_per_chapter} slides per chapter)")
                        
                        # Enhance content with Gemini (chapter by chapter)
                        logger.info("Starting content enhancement phase")
                        enhanced_content = []
                        progress_bar = st.progress(0)
                        
                        for i, chapter in enumerate(chapters):
                            chapter_num = chapter['chapter_number']
                            slide_range = f"{min(chapter['slide_numbers'])}-{max(chapter['slide_numbers'])}"
                            st.info(f"🤖 Enhancing Chapter {chapter_num} (slides {slide_range})...")
                            logger.debug(f"Enhancing content for chapter {chapter_num}")
                            
                            enhanced = converter.enhance_content_with_gemini(chapter, chapter_num)
                            enhanced_content.append(enhanced)
                            progress_bar.progress((i + 1) / len(chapters))
                        
                        logger.info("Content enhancement completed, starting eBook generation")
                        
                        # Initialize buffers
                        pdf_buffer = None
                        docx_buffer = None
                        
                        # Generate based on selected format
                        if output_format in ["PDF", "Both"]:
                            st.info("📄 Generating PDF eBook with images...")
                            
                            # Prepare header and footer text
                            header_text = custom_header if custom_header.strip() else None
                            footer_text = custom_footer if custom_footer.strip() else None
                            
                            pdf_buffer = converter.create_pdf_ebook(
                                enhanced_content, 
                                title=ebook_title,
                                author=ebook_author,
                                header_text=header_text,
                                footer_text=footer_text
                            )
                        
                        if output_format in ["DOCX", "Both"]:
                            st.info("📄 Generating DOCX eBook with images...")
                            
                            docx_buffer = converter.create_docx_ebook(
                                enhanced_content,
                                title=ebook_title,
                                author=ebook_author
                            )
                        
                        # Show success and download buttons
                        if (output_format == "PDF" and pdf_buffer) or \
                           (output_format == "DOCX" and docx_buffer) or \
                           (output_format == "Both" and pdf_buffer and docx_buffer):
                            
                            logger.info("eBook generation process completed successfully")
                            st.success(f"🎉 eBook generated successfully in {output_format} format(s)!")
                            
                            # Download buttons based on format
                            if output_format in ["PDF", "Both"] and pdf_buffer:
                                st.download_button(
                                    label="📅 Download PDF eBook",
                                    data=pdf_buffer.getvalue(),
                                    file_name=f"{ebook_title.replace(' ', '_')}.pdf",
                                    mime="application/pdf"
                                )
                            
                            if output_format in ["DOCX", "Both"] and docx_buffer:
                                st.download_button(
                                    label="📄 Download DOCX eBook",
                                    data=docx_buffer.getvalue(),
                                    file_name=f"{ebook_title.replace(' ', '_')}.docx",
                                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                                )
                            
                            # Preview section
                            with st.expander("👀 Preview Enhanced Content"):
                                for i, content in enumerate(enhanced_content[:3]):  # Show first 3 chapters
                                    chapter_num = content.get('chapter_number', i + 1)
                                    st.markdown(f"**Chapter {chapter_num}:**")
                                    
                                    # Handle both old and new content formats
                                    if isinstance(content, dict):
                                        content_text = content['text']
                                        content_images = content.get('images', [])
                                        slide_numbers = content.get('slide_numbers', [content.get('slide_number', i + 1)])
                                        
                                        # Show which slides are included in this chapter
                                        if isinstance(slide_numbers, list) and len(slide_numbers) > 1:
                                            slide_range = f"slides {min(slide_numbers)}-{max(slide_numbers)}"
                                        else:
                                            slide_range = f"slide {slide_numbers[0] if isinstance(slide_numbers, list) else slide_numbers}"
                                        
                                        st.caption(f"📄 Content from {slide_range}")
                                        st.text_area(f"Chapter {chapter_num} Content", content_text, height=150, disabled=True)
                                        
                                        if content_images:
                                            # Count images vs diagrams
                                            images = [img for img in content_images if img.get('type', 'image') == 'image']
                                            diagrams = [img for img in content_images if img.get('type', 'image') == 'diagram']
                                            
                                            if images:
                                                st.markdown(f"📸 **Images:** {len(images)} image(s) from {slide_range}")
                                            if diagrams:
                                                st.markdown(f"📈 **Diagrams/Flowcharts:** {len(diagrams)} diagram(s) from {slide_range}")
                                    else:
                                        # Fallback for old format
                                        st.text_area(f"Content {i + 1}", content, height=150, disabled=True)
                                    
                                    if i < 2:
                                        st.markdown("---")
                                
                                if len(enhanced_content) > 3:
                                    st.info(f"... and {len(enhanced_content) - 3} more sections in your eBook!")
    
    with col2:
        st.markdown("""
        <div class="feature-card">
            <h4>💡 Tips for best results</h4>
            <ul>
                <li>Use clear slide titles</li>
                <li>Include detailed content in slides</li>
                <li>Ensure text is readable</li>
                <li>Add your Gemini API key for AI enhancement</li>
            </ul>
        </div>
        """, unsafe_allow_html=True)
        
        st.markdown("""
        <div class="feature-card">
            <h4>🔑 Get Gemini API Key</h4>
            <p>Visit <a href="https://makersuite.google.com/app/apikey" target="_blank">Google AI Studio</a> to get your free Gemini API key.</p>
        </div>
        """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()
